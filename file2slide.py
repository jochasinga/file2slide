#!/usr/bin/python

from pptx import Presentation
from pptx.util import Inches
from wand.image import Image
import os, os.path


layoutMode = {
    'TITLE'                : 0, 
    'TITLE_AND_CONTENT'    : 1, 
    'SECTION_HEADER'       : 2, 
    'SEQUE'                : 2, 
    'TWO_CONTENT'          : 3, 
    'COMPARISON'           : 4, 
    'TITLE_ONLY'           : 5, 
    'BLANK'                : 6, 
    'CONTENT_WITH_CAPTION' : 7, 
    'PICTURE_WITH_CAPTION' : 8,
    ''                     : 6
}

# Crop images
def crop_image(path):
    print "Entering crop_image()"
    subfiles = os.listdir(path)

    left = int(raw_input("LEFT CROP: "))
    top = int(raw_input("TOP CROP: "))
#    right = raw_input(float("RIGHT CROP: ").strip())
#    bottom = raw_input(float("BOTTOM CROP: ").strip())

    for sf in subfiles:
        if os.path.join(path, sf).lower().endswith(('.jpg', '.png', '.jpeg', '.gif', '.tiff',)):
            print "cropping %s" % (os.path.join(path, sf))
            with Image(filename=os.path.join(path, sf)) as img:
                img.crop(left=left, top=top, width=img.width, height=img.height)
                img.save(filename=os.path.join(path, sf))
                #yield path

def pdf2image(path, *pages):
    # converting first page into JPG
    if pages:
        for page in pages:
            newpath = path + ('[%s]' %  page)
            with Image(filename=newpath) as img:
                imagepath = os.path.splitext(path)[0] + '.jpg'
                img.save(filename=imagepath)
                yield imagepath

# Filter files and images
def filter_files(path):
    files = os.listdir(path)
    for f in files:
        root = os.path.join(path, f)
        if os.path.isdir(root):
            print "Expecting FILE, got DIR!"
            if os.path.basename(root) == 'crop':
                print "Found a subdirectory named 'crop'"
                print """ 
######################## CROP IMAGES #######################\r
  + Set CROPPING for all images inside 'crop' directory.\r
  + The values are LEFT, TOP, RIGHT, and BOTTOM.\r
  + OR /images for relative path.\r
############################################################\n
                """ 
                # This doesn't run
                crop_image(root)
                '''
                if sf.lower().endswith(('.jpg', '.png', '.jpeg', '.gif', '.tiff',)):
                    crop_image(os.path.join(root, sf))

                elif sf.lower().endswith('.pdf'):
                    pdf2image(os.path.join(root, sf), 0)
                '''
        elif os.path.isfile(root):
            if root.lower().endswith(('.jpg', '.png', '.jpeg', '.gif', '.tiff',)):
                yield root
            
            elif root.lower().endswith('.pdf'):
                pdf2image(root, 0)

print """ 
#################### LOCATE DIRECTORY #######################\r
  + Locate the directory where your images are located.\r
  + For example: User/Macbook/Pictures for absolute path,\r
  + OR /images for relative path.\r
  + Make sure no subdirectories are present in the directory.\r
  + Optionally, you can drag the directory into the terminal\r
  + window after the prompt.\r
#############################################################\n
      """             
img_path = raw_input("Where is the images folder?: ").strip()

# Create a presentation file
print "Creating presentation..."
prs = Presentation()

print """
##################### CHOOSE LAYOUT STYLE ######################\r
  + Powerpoint comes with several layout styles for slides.\r
  + For example: TITLE, TITLE_WITH_CONTENT, BLANK, etc.\r
  + Type the preferred style in UPPERCASE into the next prompt\r
  + OR hit RETURN for default BLANK style.\r
"""
for mode in (k for k in layoutMode if k != ''):
    print mode

print """
################################################################\r
"""
 
layout_mode = raw_input("\nWhat's your slides layout style?: ").strip()
slide_layout = prs.slide_layouts[layoutMode[layout_mode]]

print"""
######################## SET MARGINS ###########################\r
  + Set LEFT, TOP, and RIGHT margins of your images.\r
  + Note that if RIGHT margin is set, images will be scaled\r
  + proportionally to fit. Otherwise, hit RETURN when\r
  + prompted to set margin to 0 (fit to the slide).\r
  + Margins are in inches.\r
#################################################################\n
"""
left = Inches(float( raw_input("LEFT MARGIN: ") or 0 ))
top = Inches(float( raw_input("TOP MARGIN: ") or 0 ))
width = prs.slide_width - (left + Inches(float(raw_input("RIGHT MARGIN: ") or 0)))

for path in filter_files(img_path):

    print "Creating slide..."
    slide = prs.slides.add_slide(slide_layout)

    print "Adding " + path
    pic = slide.shapes.add_picture(path, left, top, width)

print"""
##################### SAVE TO DIRECTORY ########################\r
  + CONGRATS! I finished adding images to slides alright.\r
  + Now tell me where to save your powerpoint file.\r
  + If you provide me with just a name i.e. 'test.pptx',\r
  + I will save the file to your images directory. Otherwise,\r
  + give a path like 'User/Macbook/Documents/test.pptx'\r
  + or drag the directory into this window as usual.\r
#################################################################\n
"""
save_to = raw_input("Where to save your powerpoint to?: ").strip()

# save_to = 'test.pptx'
if save_to.rpartition('/')[0] == '' and save_to.rpartition('/')[1] == '':
    if not save_to.lower().endswith('.pptx'):
        prs.save(os.path.join(img_path, save_to + '.pptx'))
        print "Your file is saved to -> " + os.path.join(img_path, save_to + '.pptx')
    else:
        prs.save(os.path.join(img_path, save_to))
        print "Your file is saved to -> " + os.path.join(img_path, save_to)
    
elif save_to.rpartition('/')[0] != '' and save_to.lower().endswith('.pptx'):
    # '/' found, look like a absolute path
    prs.save(save_to)
    print "Your file is saved to -> " + save_to

elif save_to.rpartition('/')[0] != '' and not save_to.endswith('.pptx'):
    print "Look like you have a path, but still missing the file name..."
    name = raw_input("Please type your preferred file name: ")
    name = name if name.endswith('.pptx') else (name + '.pptx')
    
    prs.save(os.path.join(save_to, name))
    print "Your file is saved to -> " + os.path.join(save_to, name)

else:
    print "There's something fishy with the file name and directory. Would you mind starting over?"




